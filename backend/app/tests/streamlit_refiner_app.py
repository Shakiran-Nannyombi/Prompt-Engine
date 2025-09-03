import uuid
import streamlit as st
import os
import sys
import chromadb
import tempfile
from pathlib import Path
import pdfplumber

# Adding the parent directory to the path so we can import from agents
sys.path.append(os.path.abspath(".."))

from dotenv import load_dotenv
from agents.refiner_agent import refiner_graph, extract_message_content, RefinerState
from langchain_core.messages import HumanMessage, AIMessage


# Load environment variables from .env file
load_dotenv(dotenv_path="../.env")

# Page configuration
st.set_page_config(
    page_title="Prompt Refiner Agent",
    layout="wide"
)

def initialize_session_state():
    if "thread_id" not in st.session_state:
        st.session_state.thread_id = str(uuid.uuid4())
    
    if "messages" not in st.session_state:
        st.session_state.messages = []
    
    if "conversation_started" not in st.session_state:
        st.session_state.conversation_started = False
    
    if "graph_state" not in st.session_state:
        st.session_state.graph_state = {}
    
    # User authentication placeholder - will be integrated with frontend auth
    if "current_user" not in st.session_state:
        st.session_state.current_user = {
            "id": "user_123",
            "name": "Demo User",
            "email": "demo@example.com"
        }
        
def main():
    st.title("Prompt Refiner Agent")
    st.markdown("**Master the art of prompt refinement with AI-powered analysis and enhancement!**")
    
    # Initializing session state
    initialize_session_state()
    
    # Sidebar for configuration
    with st.sidebar:
        
        # User info section
        st.header("User Info")
        user = st.session_state.current_user
        st.info(f"**Logged in as:** {user['name']}")
        st.caption(f"Email: {user['email']}")
        
        # Document upload section
        st.subheader("Document Upload")
        uploaded_file = st.file_uploader(
            "Upload documents for context-aware refinement",
            type=['txt', 'md', 'pdf', 'docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls', 'csv', 
                  'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'webp', 'json', 'xml', 
                  'html', 'htm', 'rtf', 'py', 'js', 'java', 'cpp', 'c', 'cs', 'php', 
                  'rb', 'go', 'rs', 'ts', 'jsx', 'tsx'],
            help="Upload documents, images, code files, or presentations to enable RAG capabilities. Supports OCR for images."
        )
        
        if uploaded_file is not None:
            # Processing the uploaded file using the existing RAG tools
            st.info(f"Processing {uploaded_file.name}...")
            
            # Create progress bar
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            success = process_uploaded_document_with_rag(uploaded_file, progress_bar, status_text)
            
            if success:
                progress_bar.progress(100)
                status_text.text("✅ Processing complete!")
                st.success(f"{uploaded_file.name} processed and indexed successfully!")
                st.info("Document is now available for context-aware refinement.")
            else:
                progress_bar.empty()
                status_text.empty()
                st.error(f"Failed to process {uploaded_file.name}. Please try again.")
        
        st.subheader("Document Status")
        rag_status = check_rag_status()
        st.info(rag_status["status"])
        
        # Clearing conversation button
        if st.button("Clear Conversation", type="secondary"):
            if 'refiner_messages' in st.session_state:
                del st.session_state.refiner_messages
            st.rerun()
    
    # Main chat interface
    # Initializing session state
    if 'refiner_messages' not in st.session_state:
        st.session_state.refiner_messages = []
    
    # Display conversation history
    for message in st.session_state.refiner_messages:
        role, content = extract_message_content(message)
        
        if role == "user":
            with st.chat_message("user"):
                st.markdown(content)
        elif role == "assistant":
            with st.chat_message("assistant"):
                st.markdown(content)
        elif role == "system":
            with st.chat_message("assistant"):
                st.markdown(f"**System:** {content}")
    
    # Suggestions and controls
    col1, col2, col3 = st.columns([2, 1, 2])
    with col2:
        if st.button("Ask for Suggestions", type="secondary", use_container_width=True):
            # Create progress bar for suggestions
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            status_text.text("🤔 Generating suggestions...")
            progress_bar.progress(30)
            
            suggestions_message = get_refinement_suggestions("", st.session_state.thread_id)
            
            progress_bar.progress(70)
            status_text.text("📝 Formatting suggestions...")
            
            if suggestions_message:
                st.session_state.refiner_messages.append(suggestions_message)
                
                progress_bar.progress(100)
                status_text.text("✅ Suggestions ready!")
                
                # Show the suggestions immediately
                with st.chat_message("assistant"):
                    st.markdown(suggestions_message.content)
                    
                # Clean up progress indicators
                progress_bar.empty()
                status_text.empty()
                st.rerun()

    with col3:
        if st.button("New Session", type="secondary", use_container_width=True):
            clear_conversation()    # Chat input
    if prompt := st.chat_input("Enter your prompt for refinement...", key="user_input"):
        # Checking if this thread already has a refinement
        has_refinement = any(
            hasattr(msg, 'content') and ("Refined Prompt:" in msg.content or "Final Report:" in msg.content) 
            for msg in st.session_state.refiner_messages
        )
        
        # Checking if this is a follow-up question about the existing refinement
        is_follow_up = any(
            keyword in prompt.lower() for keyword in 
            ["refined", "prompt", "framework", "category", "analysis", "explain", "why", "how", "what", "help", "suggestions", "recommendations"]
        )
        
        # Checking if this is a greeting (should always be allowed)
        is_greeting = any(
            keyword in prompt.lower() for keyword in 
            ["hello", "hi", "hey", "greetings", "good morning", "good afternoon", "good evening", "how are you", "what's up"]
        )
        
        if has_refinement and not is_follow_up and not is_greeting:
            # User is trying to refine a new prompt in the same thread
            with st.chat_message("user"):
                st.markdown(prompt)
            
            with st.chat_message("assistant"):
                st.warning("""
⚠️ **New Prompt Detected**

I see you want to refine a new prompt, but this chat thread already contains a previous refinement.

**To refine a new prompt:**
1. Click "New Session" to start a fresh chat
2. Or ask me questions about your current refinement

**To continue with this refinement:**
Ask me about the framework used, analysis, or recommendations

This helps maintain clear conversation context and prevents confusion between different refinements.
                """)
        else:
            # Add user message to session state first
            user_message = HumanMessage(content=prompt)
            st.session_state.refiner_messages.append(user_message)
            
            # Create progress bar for prompt processing
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            status_text.text("🤖 Processing your prompt...")
            progress_bar.progress(20)
            
            # Processing the message and getting the response
            response_message = process_refinement(prompt, st.session_state.thread_id, progress_bar, status_text)
            
            progress_bar.progress(90)
            status_text.text("✅ Processing complete!")
            
            # Adding response to session state if we got one
            if response_message:
                st.session_state.refiner_messages.append(response_message)
            
            progress_bar.progress(100)
            
            # Clean up progress indicators
            progress_bar.empty()
            status_text.empty()
            
            # Rerun to update the display with new messages
            st.rerun()

def clear_conversation():
    """Clear the conversation and start a new session"""
    if 'refiner_messages' in st.session_state:
        del st.session_state.refiner_messages
    if 'thread_id' in st.session_state:
        st.session_state.thread_id = str(uuid.uuid4())
    
    # Initialize fresh message list
    st.session_state.refiner_messages = []
    
    # Show success message
    st.success(" New session started! You can now refine a new prompt.")
    st.rerun()

# Processing the user's prompt through the refiner agent
def process_refinement(user_input: str, thread_id: str, progress_bar=None, status_text=None):
    """Process user input through the refiner agent"""
    try:
        # Step 1: Prepare input
        if progress_bar:
            progress_bar.progress(40)
        if status_text:
            status_text.text("📝 Preparing prompt for analysis...")
        
        # Preparing input for the refiner graph
        user_message = HumanMessage(content=user_input)
        graph_input = {
            "messages": [user_message],
            "text": user_input,
            "refined_text": "",
            "prompt_category": "",
            "framework_used": "",
            "analysis": "",
            "original_prompt": user_input,
            "has_uploaded_documents": False,
            "document_context": "",
            "is_greeting": False,
            "is_help_request": False,
            "is_follow_up": False,
            "continue_to_refinement": False
        }
        
        # Check if documents are available for RAG
        if os.path.exists("./chroma_db"):
            graph_input["has_uploaded_documents"] = True
            graph_input["document_context"] = "Document context available for enhanced refinement"
            if progress_bar:
                progress_bar.progress(50)
            if status_text:
                status_text.text("📚 Loading document context...")

        # Step 2: Process with AI agents
        if progress_bar:
            progress_bar.progress(60)
        if status_text:
            status_text.text("🤖 Processing with AI agents...")

        # Invoke the refiner graph
        result = refiner_graph.invoke(
            graph_input,
            config={"configurable": {"thread_id": thread_id}}
        )
        
        # Step 3: Extract results
        if progress_bar:
            progress_bar.progress(80)
        if status_text:
            status_text.text("Extracting results...")
        
        # Extract the final message and return it
        if result.get("messages"):
            final_message = result["messages"][-1]
            return final_message
        else:
            # Create an error message if no response
            error_msg = "No response received from the refiner agent. Please try again."
            if status_text:
                status_text.text(f"❌ {error_msg}")
            else:
                st.error(error_msg)
            return AIMessage(content=error_msg)
            
    except Exception as e:
        # Create an error message if there's an exception
        error_msg = f"Error processing refinement: {str(e)}"
        if status_text:
            status_text.text(f"❌ {error_msg}")
        else:
            st.error(error_msg)
        return AIMessage(content=error_msg)

# Getting suggestions for prompt refinement without full processing
def get_refinement_suggestions(user_input: str, thread_id: str):
    try:
        # Create a suggestion response - handle case where session state might not be available
        try:
            user_name = st.session_state.current_user['name']
        except (AttributeError, KeyError):
            user_name = "User"  
        
        if user_input.strip():
            # Specific suggestions for the provided input
            suggestion_content = f"""
Hi {user_name}! Here are some suggestions to improve your prompt:

**Current Prompt:** {user_input}

**Suggestions:**
1. **Be Specific**: Add concrete details about what you want to achieve
2. **Define Context**: Specify the background, audience, and constraints
3. **Set Expectations**: Clearly state the desired output format and quality
4. **Include Examples**: Provide sample inputs/outputs if possible
5. **Specify Constraints**: Mention any limitations or requirements
            """
        else:
            # General suggestions for prompt engineering
            suggestion_content = f"""
Hi {user_name}! Here are some general prompt engineering suggestions:

**Key Prompt Engineering Techniques:**

1. **Zero-Shot Prompting**: Direct instructions without examples
2. **Few-Shot Prompting**: Include 1-3 examples in your prompt
3. **Chain-of-Thought**: Ask the AI to show its reasoning step-by-step
4. **Role-Based Prompting**: Assign a specific role or persona to the AI
5. **Template-Based**: Use structured formats with placeholders

**Best Practices:**
- Be clear and specific about what you want
- Define the context and background
- Specify the desired output format
- Include examples when helpful
- Set constraints and limitations
- Use action words and clear instructions

Ready to refine a prompt? Just type it in the chat below! 
            """

        suggestion_message = AIMessage(content=suggestion_content)
        return suggestion_message
        
    except Exception as e:
        return AIMessage(content=f"Error generating suggestions: {str(e)}")
    
    # Fallback return
    return None

def process_uploaded_document_with_rag(uploaded_file, progress_bar=None, status_text=None):
    try:
        # Step 1: Extract text content based on file type
        if progress_bar:
            progress_bar.progress(10)
        if status_text:
            status_text.text("Extracting text from document...")
            
        text_content = extract_text_from_uploaded_file(uploaded_file, progress_bar, status_text)
        
        if not text_content:
            if status_text:
                status_text.text("❌ Could not extract text from the uploaded file.")
            return False
        
        if len(text_content.strip()) < 10:
            if status_text:
                status_text.text("⚠️ The extracted text is too short to process meaningfully.")
            return False
        
        # Step 2: Process with RAG tools
        if progress_bar:
            progress_bar.progress(60)
        if status_text:
            status_text.text("🧠 Processing with AI and creating embeddings...")
        
        # Using the existing RAG tool to process the file
        from agents.tools.refinement_tools import process_uploaded_file
        
        result = process_uploaded_file(text_content)
        
        # Step 3: Finalize
        if progress_bar:
            progress_bar.progress(90)
        if status_text:
            status_text.text("💾 Storing in vector database...")
        
        # Checking if processing was successful
        if "successfully" in result.lower():
            return True
        else:
            if status_text:
                status_text.text(f"❌ Processing failed: {result}")
            return False
            
    except Exception as e:
        if status_text:
            status_text.text(f"❌ Error processing document: {str(e)}")
        return False

def extract_text_from_uploaded_file(uploaded_file, progress_bar=None, status_text=None):
    try:
        # Creating temporary file to save uploaded content
        file_extension = Path(uploaded_file.name).suffix.lower()
        
        if progress_bar:
            progress_bar.progress(20)
        if status_text:
            status_text.text(f"📁 Preparing {file_extension} file for processing...")
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name
        
        try:
            # Extracting text based on file type
            if progress_bar:
                progress_bar.progress(30)
            if status_text:
                status_text.text(f"📖 Extracting text from {file_extension} file...")
                
            text_content = ""
            
            # Text files
            if file_extension in ['.txt', '.md', '.rtf']:
                with open(tmp_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text_content = f.read()
            
            # PDF files
            elif file_extension == '.pdf':
                text_content = extract_pdf_text(tmp_file_path)
            
            # Word documents
            elif file_extension in ['.docx', '.doc']:
                text_content = extract_word_text(tmp_file_path)
            
            # PowerPoint presentations
            elif file_extension in ['.pptx', '.ppt']:
                text_content = extract_powerpoint_text(tmp_file_path)
            
            # Excel spreadsheets
            elif file_extension in ['.xlsx', '.xls', '.csv']:
                text_content = extract_excel_text(tmp_file_path)
            
            # Image files (OCR)
            elif file_extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']:
                text_content = extract_image_text(tmp_file_path)
            
            # JSON files
            elif file_extension == '.json':
                import json
                with open(tmp_file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    text_content = json.dumps(data, indent=2)
            
            # XML files
            elif file_extension == '.xml':
                with open(tmp_file_path, 'r', encoding='utf-8') as f:
                    text_content = f.read()
            
            # HTML files
            elif file_extension in ['.html', '.htm']:
                text_content = extract_html_text(tmp_file_path)
            
            # Code files
            elif file_extension in ['.py', '.js', '.java', '.cpp', '.c', '.cs', '.php', '.rb', '.go', '.rs', '.ts', '.jsx', '.tsx']:
                with open(tmp_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text_content = f.read()
            
            # Other text-based files
            else:
                # Try to read them as text files
                try:
                    with open(tmp_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text_content = f.read()
                except:
                    st.warning(f"Unsupported file format: {file_extension}. Trying OCR...")
                    # As a last resort, try OCR if it's an image-like file
                    try:
                        text_content = extract_image_text(tmp_file_path)
                    except:
                        return None
            
            # Final progress update
            if progress_bar:
                progress_bar.progress(50)
            if status_text:
                status_text.text(f"✅ Text extracted successfully from {file_extension} file!")
                
            return text_content
            
        finally:
            # Cleaning up temporary file
            Path(tmp_file_path).unlink(missing_ok=True)
            
    except Exception as e:
        if status_text:
            status_text.text(f"❌ Error extracting text from file: {str(e)}")
        else:
            st.error(f"Error extracting text from file: {str(e)}")
        return None

def extract_pdf_text(file_path):
    """Extract text from PDF files"""
    try:
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text
    except ImportError:
        # Fallback to PyPDF2
        try:
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except:
            return None
    except Exception as e:
        st.warning(f"Error reading PDF: {str(e)}")
        return None

def extract_word_text(file_path):
    """Extract text from Word documents"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        
        return text
    except ImportError:
        st.warning("python-docx not available. Word document processing disabled.")
        return None
    except Exception as e:
        st.warning(f"Error reading Word document: {str(e)}")
        return None

def extract_powerpoint_text(file_path):
    """Extract text from PowerPoint presentations"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        
        return text
    except ImportError:
        st.warning("python-pptx not available. PowerPoint processing disabled.")
        return None
    except Exception as e:
        st.warning(f"Error reading PowerPoint: {str(e)}")
        return None

def extract_excel_text(file_path):
    """Extract text from Excel spreadsheets"""
    try:
        import openpyxl
        import pandas as pd
        
        if file_path.endswith('.csv'):
            # Handle CSV files
            df = pd.read_csv(file_path)
            return df.to_string()
        else:
            # Handle Excel files
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text += " | ".join(row_text) + "\n"
                text += "\n"
            
            return text
    except ImportError:
        st.warning("openpyxl not available. Excel processing disabled.")
        return None
    except Exception as e:
        st.warning(f"Error reading Excel file: {str(e)}")
        return None

def extract_image_text(file_path):
    """Extract text from images using OCR"""
    try:
        import pytesseract
        from PIL import Image
        
        # Open and process the image
        image = Image.open(file_path)
        
        # Convert to RGB if necessary
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Extract text using OCR
        text = pytesseract.image_to_string(image)
        
        if not text.strip():
            st.info("No text found in the image.")
            return None
            
        return text
    except ImportError:
        st.warning("pytesseract not available. Image OCR processing disabled. Please install tesseract-ocr system package.")
        return None
    except Exception as e:
        st.warning(f"Error processing image with OCR: {str(e)}")
        return None

def extract_html_text(file_path):
    """Extract text from HTML files"""
    try:
        from bs4 import BeautifulSoup
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        text = soup.get_text()
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return text
    except ImportError:
        # Fallback: read as plain text
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        st.warning(f"Error reading HTML file: {str(e)}")
        return None

def check_rag_status():
    try:
        import os
        chroma_path = "./chroma_db"
        
        if os.path.exists(chroma_path):
            # Checking if ChromaDB folder exists with documents
            try:
                client = chromadb.PersistentClient(path=chroma_path)
                collections = client.list_collections()
                
                if collections:
                    collection = client.get_collection("uploads_collection")
                    count = collection.count()
                    return {
                        "has_documents": count > 0,
                        "document_count": count,
                        "status": f"{count} documents indexed" if count > 0 else "📭 No documents indexed"
                    }
                else:
                    return {
                        "has_documents": False,
                        "document_count": 0,
                        "status": "No documents indexed"
                    }
            except Exception:
                return {
                    "has_documents": False,
                    "document_count": 0,
                    "status": "ChromaDB folder exists but no documents"
                }
        else:
            return {
                "has_documents": False,
                "document_count": 0,
                "status": "No documents uploaded yet"
            }
    except Exception as e:
        return {
            "has_documents": False,
            "document_count": 0,
            "status": f"Error checking status: {str(e)}"
        }

if __name__ == "__main__":
    main()
